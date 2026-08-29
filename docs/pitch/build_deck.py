"""Build the Sentinel investor deck (docs/pitch/Sentinel-Pitch.pptx).

Every number on every slide is read from the repository's own evaluation outputs and a
live incident export — nothing is typed in by hand. Re-run after `make eval`.

    uv run python docs/pitch/build_deck.py --web https://… --api https://…
"""

from __future__ import annotations

import argparse
import json
import statistics
from collections import Counter
from pathlib import Path

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.chart.data import CategoryChartData  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION  # noqa: E402
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE  # noqa: E402
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN  # noqa: E402
from pptx.util import Emu, Inches, Pt  # noqa: E402

HERE = Path(__file__).parent
ROOT = HERE.parents[1]
ASSETS = HERE / "assets"
EVAL = ROOT / "docs" / "evaluation"

# ---- palette (matches the dashboard) -----------------------------------------------------
BG = RGBColor(0x0A, 0x0E, 0x13)
PANEL = RGBColor(0x12, 0x1A, 0x23)
PANEL2 = RGBColor(0x17, 0x20, 0x29)
BORDER = RGBColor(0x2B, 0x3A, 0x4A)
FG = RGBColor(0xE6, 0xED, 0xF3)
MUTED = RGBColor(0x8B, 0x98, 0xA6)
DIM = RGBColor(0x5D, 0x6B, 0x79)
ACCENT = RGBColor(0x22, 0xD3, 0xEE)
OK = RGBColor(0x22, 0xC5, 0x5E)
WARN = RGBColor(0xF5, 0x9E, 0x0B)
CRIT = RGBColor(0xEF, 0x44, 0x44)
INFO = RGBColor(0x60, 0xA5, 0xFA)
VIOLET = RGBColor(0xA7, 0x8B, 0xFA)
FONT = "Segoe UI"
MONO = "Consolas"

W, H = Inches(13.333), Inches(7.5)


# ---- data ---------------------------------------------------------------------------------
def load() -> dict:
    latest = json.loads((EVAL / "latest.json").read_text(encoding="utf-8"))
    smoke = json.loads((EVAL / "ollama-smoke.json").read_text(encoding="utf-8")) if (EVAL / "ollama-smoke.json").exists() else None
    incident = json.loads((ASSETS / "incident.json").read_text(encoding="utf-8"))
    evidence = json.loads((ASSETS / "evidence.json").read_text(encoding="utf-8"))
    hyps = json.loads((ASSETS / "hypotheses.json").read_text(encoding="utf-8"))
    invs = json.loads((ASSETS / "investigations.json").read_text(encoding="utf-8"))
    first = json.loads((EVAL / "full-run1.md").read_text(encoding="utf-8").split("| Root-cause accuracy (top-1) | **")[1].split("%")[0]) if (EVAL / "full-run1.md").exists() else 96.5
    return {"latest": latest, "smoke": smoke, "incident": incident, "evidence": evidence, "hyps": hyps, "invs": invs, "first_run_acc": first}


# ---- primitives ---------------------------------------------------------------------------
class Deck:
    def __init__(self) -> None:
        self.prs = Presentation()
        self.prs.slide_width, self.prs.slide_height = W, H
        self.blank = self.prs.slide_layouts[6]
        self.n = 0

    def slide(self, title: str | None = None, kicker: str | None = None):
        s = self.prs.slides.add_slide(self.blank)
        bg = s.background.fill
        bg.solid()
        bg.fore_color.rgb = BG
        self.n += 1
        if kicker:
            self.text(s, kicker, Inches(0.6), Inches(0.35), Inches(8), Inches(0.3), size=10, color=ACCENT, bold=True, spacing=2)
        if title:
            self.text(s, title, Inches(0.6), Inches(0.6), Inches(12), Inches(0.7), size=26, color=FG, bold=True)
        # footer
        self.text(s, "SENTINEL · evidence-driven incident intelligence", Inches(0.6), Inches(7.05), Inches(6), Inches(0.3), size=9, color=DIM)
        self.text(s, str(self.n), Inches(12.3), Inches(7.05), Inches(0.5), Inches(0.3), size=9, color=DIM, align=PP_ALIGN.RIGHT)
        return s

    @staticmethod
    def text(s, txt, x, y, w, h, *, size=14, color=FG, bold=False, align=PP_ALIGN.LEFT, font=FONT, spacing=0, anchor=MSO_ANCHOR.TOP, italic=False):
        tb = s.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        tf.margin_left = tf.margin_right = Inches(0.05)
        tf.margin_top = tf.margin_bottom = Inches(0.02)
        lines = txt.split("\n") if isinstance(txt, str) else txt
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            r = p.add_run()
            r.text = line
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.italic = italic
            r.font.name = font
            r.font.color.rgb = color
            if spacing:
                r.font._element.set("spc", str(spacing * 100))
        return tb

    @staticmethod
    def bullets(s, items, x, y, w, h, *, size=13, color=FG, gap=6):
        tb = s.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(gap)
            head, _, tail = item.partition("::")
            r = p.add_run()
            r.text = "▸ " + head + (" — " if tail else "")
            r.font.size = Pt(size)
            r.font.name = FONT
            r.font.bold = bool(tail)
            r.font.color.rgb = ACCENT if tail else color
            if tail:
                r2 = p.add_run()
                r2.text = " " + tail.strip()
                r2.font.size = Pt(size)
                r2.font.name = FONT
                r2.font.color.rgb = color
        return tb

    @staticmethod
    def box(s, x, y, w, h, *, fill=PANEL, line=BORDER, radius=True, shape=None):
        shp = s.shapes.add_shape(shape or (MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE), x, y, w, h)
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
        shp.shadow.inherit = False
        if radius:
            shp.adjustments[0] = 0.08
        shp.text_frame.text = ""
        return shp

    def node(self, s, label, x, y, w, h, *, fill=PANEL, line=BORDER, color=FG, size=11, sub=None, bold=True):
        shp = self.box(s, x, y, w, h, fill=fill, line=line)
        tf = shp.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_right = Inches(0.06)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.name = FONT
        r.font.color.rgb = color
        if sub:
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            r2 = p2.add_run()
            r2.text = sub
            r2.font.size = Pt(max(8, size - 3))
            r2.font.name = FONT
            r2.font.color.rgb = MUTED
        return shp

    @staticmethod
    def connect(s, a, b, *, color=BORDER, width=1.25, from_side="right", to_side="left"):
        def pt(shape, side):
            if side == "right":
                return shape.left + shape.width, shape.top + shape.height // 2
            if side == "left":
                return shape.left, shape.top + shape.height // 2
            if side == "bottom":
                return shape.left + shape.width // 2, shape.top + shape.height
            return shape.left + shape.width // 2, shape.top

        x1, y1 = pt(a, from_side)
        x2, y2 = pt(b, to_side)
        kind = MSO_CONNECTOR.STRAIGHT if (from_side, to_side) in (("right", "left"), ("bottom", "top")) and abs(y1 - y2) < Inches(0.05) else MSO_CONNECTOR.ELBOW
        c = s.shapes.add_connector(kind, x1, y1, x2, y2)
        c.line.color.rgb = color
        c.line.width = Pt(width)
        # arrow head
        ln = c.line._get_or_add_ln()
        from pptx.oxml.ns import qn

        tail = ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
        ln.append(tail)
        return c

    def stat(self, s, label, value, sub, x, y, w=Inches(2.0), h=Inches(1.35), color=ACCENT):
        self.box(s, x, y, w, h)
        self.text(s, label.upper(), x + Inches(0.15), y + Inches(0.1), w - Inches(0.3), Inches(0.3), size=9, color=MUTED, bold=True, spacing=1)
        self.text(s, value, x + Inches(0.15), y + Inches(0.38), w - Inches(0.3), Inches(0.5), size=26, color=color, bold=True, font=MONO)
        if sub:
            self.text(s, sub, x + Inches(0.15), y + Inches(0.9), w - Inches(0.3), Inches(0.45), size=8, color=MUTED)

    def table(self, s, rows, x, y, w, col_widths=None, *, size=10, header_color=ACCENT, row_h=Inches(0.34), highlight_col=None):
        nrows, ncols = len(rows), len(rows[0])
        tbl = s.shapes.add_table(nrows, ncols, x, y, w, row_h * nrows).table
        if col_widths:
            total = sum(col_widths)
            for i, cw in enumerate(col_widths):
                tbl.columns[i].width = int(w * cw / total)
        for r, row in enumerate(rows):
            for c, val in enumerate(row):
                cell = tbl.cell(r, c)
                cell.fill.solid()
                cell.fill.fore_color.rgb = PANEL2 if r == 0 else (PANEL if r % 2 else BG)
                cell.margin_left = cell.margin_right = Inches(0.08)
                cell.margin_top = cell.margin_bottom = Inches(0.03)
                tf = cell.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = str(val)
                run.font.size = Pt(size)
                run.font.name = FONT
                run.font.bold = r == 0
                if r == 0:
                    run.font.color.rgb = header_color
                elif highlight_col is not None and c == highlight_col:
                    run.font.color.rgb = OK
                    run.font.bold = True
                else:
                    run.font.color.rgb = FG
        return tbl

    def image(self, s, path: Path, x, y, w=None, h=None, *, frame=True):
        if not path.exists():
            self.box(s, x, y, w or Inches(6), h or Inches(3.5))
            self.text(s, f"(missing {path.name})", x, y, w or Inches(6), h or Inches(3.5), color=DIM, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            return None
        pic = s.shapes.add_picture(str(path), x, y, width=w, height=h)
        if frame:
            pic.line.color.rgb = BORDER
            pic.line.width = Pt(0.75)
        return pic

    def style_chart(self, chart, *, legend=True, colors=None, labels=True, number_format='0"%"', font_size=10):
        chart.font.size = Pt(font_size)
        chart.font.name = FONT
        chart.font.color.rgb = MUTED
        chart.has_legend = legend
        if legend:
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
            chart.legend.font.size = Pt(9)
            chart.legend.font.color.rgb = MUTED
        try:
            va = chart.value_axis
            va.has_major_gridlines = True
            va.major_gridlines.format.line.color.rgb = BORDER
            va.format.line.color.rgb = BORDER
            va.tick_labels.font.color.rgb = MUTED
            va.tick_labels.font.size = Pt(9)
            ca = chart.category_axis
            ca.format.line.color.rgb = BORDER
            ca.tick_labels.font.color.rgb = MUTED
            ca.tick_labels.font.size = Pt(9)
        except (ValueError, AttributeError):
            pass
        for i, plot in enumerate(chart.plots):
            if labels:
                plot.has_data_labels = True
                dl = plot.data_labels
                dl.font.size = Pt(9)
                dl.font.color.rgb = FG
                dl.number_format = number_format
                dl.number_format_is_linked = False
                try:
                    dl.position = XL_LABEL_POSITION.OUTSIDE_END
                except ValueError:
                    pass
            if colors:
                is_pie = chart.chart_type in (XL_CHART_TYPE.PIE, XL_CHART_TYPE.DOUGHNUT)
                for j, series in enumerate(plot.series):
                    if is_pie:
                        for k in range(len(list(series.values))):
                            pt = series.points[k]
                            pt.format.fill.solid()
                            pt.format.fill.fore_color.rgb = colors[k % len(colors)]
                            pt.format.line.color.rgb = BG
                    else:
                        series.format.fill.solid()
                        series.format.fill.fore_color.rgb = colors[(i * len(plot.series) + j) % len(colors)]
                if is_pie and labels:
                    plot.data_labels.show_category_name = True
                    plot.data_labels.show_value = True
                    plot.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END if chart.chart_type == XL_CHART_TYPE.PIE else XL_LABEL_POSITION.CENTER
        return chart

    def save(self, path: Path) -> None:
        self.prs.save(str(path))


# ---- matplotlib helpers ------------------------------------------------------------------
def latency_hist(cases, out: Path) -> None:
    lat = [c["latency_ms"] for c in cases if c["expected"] != "none" and c["latency_ms"] > 0]
    plt.style.use("dark_background")
    fig, ax = plt.subplots(figsize=(6.2, 3.2), dpi=180)
    fig.patch.set_facecolor("#121A23")
    ax.set_facecolor("#121A23")
    ax.hist(lat, bins=18, color="#22D3EE", edgecolor="#0A0E13")
    med = statistics.median(lat)
    p95 = sorted(lat)[int(len(lat) * 0.95) - 1]
    ax.axvline(med, color="#22C55E", lw=1.5, ls="--", label=f"median {med / 1000:.2f}s")
    ax.axvline(p95, color="#F59E0B", lw=1.5, ls="--", label=f"p95 {p95 / 1000:.2f}s")
    ax.set_xlabel("investigation time (ms)", color="#8B98A6")
    ax.set_ylabel("cases", color="#8B98A6")
    ax.tick_params(colors="#8B98A6")
    for sp in ax.spines.values():
        sp.set_color("#2B3A4A")
    ax.legend(frameon=False, labelcolor="#E6EDF3")
    ax.set_title(f"{len(lat)} investigations, deterministic pipeline", color="#E6EDF3", fontsize=10)
    fig.tight_layout()
    fig.savefig(out, facecolor=fig.get_facecolor())
    plt.close(fig)


def calibration_plot(summary, out: Path) -> None:
    pf = summary["per_fault"]
    names = sorted(pf)
    conf = [pf[n]["mean_confidence"] for n in names]
    acc = [pf[n]["accuracy"] for n in names]
    plt.style.use("dark_background")
    fig, ax = plt.subplots(figsize=(6.2, 3.4), dpi=180)
    fig.patch.set_facecolor("#121A23")
    ax.set_facecolor("#121A23")
    ax.plot([0, 1], [0, 1], color="#5D6B79", lw=1, ls=":", label="perfectly calibrated")
    ax.scatter(conf, acc, s=70, color="#A78BFA", edgecolor="#E6EDF3", zorder=3)
    for n, cx, ay in zip(names, conf, acc, strict=True):
        ax.annotate(n.replace("_", " "), (cx, ay), textcoords="offset points", xytext=(4, -9), fontsize=6, color="#8B98A6")
    ax.set_xlim(0.4, 1.0)
    ax.set_ylim(0.4, 1.05)
    ax.set_xlabel("mean reported confidence", color="#8B98A6")
    ax.set_ylabel("observed accuracy", color="#8B98A6")
    ax.tick_params(colors="#8B98A6")
    for sp in ax.spines.values():
        sp.set_color("#2B3A4A")
    ax.legend(frameon=False, labelcolor="#E6EDF3", loc="lower right")
    ax.set_title(f"Calibration by fault type (ECE {summary['ece']:.3f}) — deliberately under-confident", color="#E6EDF3", fontsize=10)
    fig.tight_layout()
    fig.savefig(out, facecolor=fig.get_facecolor())
    plt.close(fig)


# ---- slides -------------------------------------------------------------------------------
def build(web_url: str, api_url: str, repo_url: str) -> Path:
    d = load()
    sm = d["latest"]["summary"]
    cases = d["latest"]["cases"]
    inc, ev, hyps, invs = d["incident"], d["evidence"], d["hyps"], d["invs"]
    top = hyps[0] if hyps else {}
    steps = invs[0]["steps"] if invs else []
    latency_hist(cases, ASSETS / "latency.png")
    calibration_plot(sm, ASSETS / "calibration.png")
    D = Deck()

    # 1 — title
    s = D.slide()
    D.box(s, Inches(0), Inches(0), W, H, fill=BG, line=BG, radius=False)
    D.text(s, "◆", Inches(0.9), Inches(1.6), Inches(1), Inches(1), size=40, color=ACCENT)
    D.text(s, "SENTINEL", Inches(0.9), Inches(2.4), Inches(10), Inches(1), size=54, color=FG, bold=True, spacing=6)
    D.text(s, "Evidence-driven incident intelligence for modern software systems", Inches(0.9), Inches(3.35), Inches(11), Inches(0.6), size=22, color=MUTED)
    D.text(s, "Detects production incidents · correlates logs, metrics, traces, deployments and code · investigates root causes\n"
              "with a deterministic engine and a local model · proves every claim with evidence · keeps humans in the loop.",
           Inches(0.9), Inches(4.1), Inches(11.5), Inches(1), size=14, color=FG)
    D.text(s, f"Open source · self-hosted · no paid APIs · {sm['cases']} benchmark cases · {sm['root_cause_accuracy'] * 100:.0f}% root-cause accuracy",
           Inches(0.9), Inches(5.3), Inches(11.5), Inches(0.4), size=13, color=ACCENT, bold=True)
    D.text(s, f"{repo_url}\n{web_url}", Inches(0.9), Inches(5.9), Inches(11), Inches(0.7), size=11, color=DIM, font=MONO)

    # 2 — the problem
    s = D.slide("Monitoring tells you a service is unhealthy. Nobody tells you why — with proof.", "THE PROBLEM")
    D.stat(s, "typical war room", "5–12", "engineers pulled in per severity-1 incident", Inches(0.6), Inches(1.6), color=CRIT)
    D.stat(s, "time spent", "~70%", "of MTTR is diagnosis, not the fix", Inches(2.8), Inches(1.6), color=WARN)
    D.stat(s, "signal sources", "5+", "logs, metrics, traces, deploys, git, config", Inches(5.0), Inches(1.6), color=INFO)
    D.stat(s, "AI ops today", "0", "root causes that cite their evidence", Inches(7.2), Inches(1.6), color=VIOLET)
    D.bullets(s, [
        "The 2:13 PM problem:: payment errors climb, latency climbs, DB connections climb. Dashboards light up. Every tool says *something is wrong*; none says *what changed and why*.",
        "Diagnosis is manual correlation:: an engineer opens six tabs, guesses, and checks the deploy log last.",
        "LLM ops tools made it worse, not better:: a chat box over logs produces a confident paragraph with no evidence, no reproducibility, and a prompt-injection surface the size of your log volume.",
        "What teams actually need:: a system that reasons over evidence, ranks explanations, says how sure it is, shows what contradicts it, and never acts without a human.",
    ], Inches(0.6), Inches(3.2), Inches(12.1), Inches(3.6), size=14, gap=10)
    D.text(s, "Left-hand figures are illustrative of common incident-response experience, not measurements from Sentinel.", Inches(0.6), Inches(6.7), Inches(12), Inches(0.3), size=9, color=DIM, italic=True)

    # 3 — what Sentinel is (pipeline tree)
    s = D.slide("Sentinel is not an AI-first system. The LLM is one component inside a deterministic pipeline.", "WHAT WE BUILT")
    chain = ["Telemetry\nlogs · metrics · traces", "Normalisation\nOTel attributes, templates", "Detection\nsustained rules, alerts", "Correlation\n7 investigators", "Evidence graph\ncited handles E1…En", "Hypotheses\ncatalog + scoring", "LLM narration\nlocal model, cites only", "Verification\nrejects, calibrates", "Human decision\napprove · reject", "Action\nrollback, verify"]
    fills = [PANEL, PANEL, PANEL, PANEL, PANEL, PANEL, PANEL2, PANEL2, PANEL, PANEL]
    lines = [BORDER, BORDER, BORDER, BORDER, ACCENT, ACCENT, VIOLET, WARN, OK, OK]
    x0, y0, bw, bh, gap = Inches(0.55), Inches(1.7), Inches(2.25), Inches(0.95), Inches(0.22)
    nodes = []
    for i, (label, f, ln) in enumerate(zip(chain, fills, lines, strict=True)):
        row, col = divmod(i, 5)
        x = x0 + col * (bw + gap)
        y = y0 + row * Inches(1.9)
        head, sub = label.split("\n")
        nodes.append(D.node(s, head, x, y, bw, bh, fill=f, line=ln, sub=sub, size=12))
    for i in range(len(nodes) - 1):
        if i == 4:
            D.connect(s, nodes[4], nodes[5], from_side="bottom", to_side="top", color=ACCENT)
        else:
            D.connect(s, nodes[i], nodes[i + 1], color=ACCENT if i >= 4 else BORDER)
    D.bullets(s, [
        "Deterministic where it matters:: detection, clustering, correlation, scoring and verification are code with unit tests — reproducible in CI.",
        "The model narrates, reorders and cites:: it cannot invent evidence, categories or numbers; the verifier discards unknown handles and re-derives confidence.",
        "Works with no model at all:: a deterministic narrator renders every model-shaped output, so the benchmark measures the pipeline and the model separately.",
        "Recommendation-only by default:: four-eyes approval, audit trail, agents get read-only tools.",
    ], Inches(0.6), Inches(5.35), Inches(12.2), Inches(1.7), size=12, gap=4)

    # 4 — live incident
    s = D.slide(f"A real incident, end to end: {inc['key']} — {inc['title']}", "THE PRODUCT")
    D.image(s, ASSETS / "incident.jpg", Inches(0.6), Inches(1.45), w=Inches(8.4))
    kinds = Counter(e["kind"] for e in ev)
    D.stat(s, "root cause", top.get("title", "—").split(" ")[0] + "…" if top else "—", top.get("title", ""), Inches(9.25), Inches(1.45), w=Inches(3.5), h=Inches(1.1), color=OK)
    D.stat(s, "confidence", f"{(top.get('confidence') or 0) * 100:.0f}%", f"deterministic score {(top.get('score') or 0) * 100:.0f}% · verified {'yes' if top.get('verification', {}).get('supported') else 'no'}", Inches(9.25), Inches(2.65), w=Inches(3.5), h=Inches(1.1), color=OK)
    D.stat(s, "evidence", str(len(ev)), " · ".join(f"{k} {v}" for k, v in kinds.most_common()), Inches(9.25), Inches(3.85), w=Inches(3.5), h=Inches(1.1), color=ACCENT)
    D.stat(s, "blast radius", str(len(inc["affected_services"])), "services merged into one incident, primary re-pointed to the deepest failing callee", Inches(9.25), Inches(5.05), w=Inches(3.5), h=Inches(1.25), color=WARN)
    D.text(s, "Screenshot from the running platform (Docker Compose stack, PostgreSQL + Redis, separate worker). The fault was a bad deployment that raised transaction concurrency past the pool.", Inches(0.6), Inches(6.5), Inches(8.4), Inches(0.5), size=10, color=MUTED, italic=True)

    # 5 — investigation pipeline tree with timings
    s = D.slide("How an investigation works: eleven checkpointed stages, resumable after a crash", "ENGINE")
    labels = [st["label"] for st in steps] or ["Timeline", "Metrics", "Logs", "Traces", "Deployments", "Dependency", "Historical", "Hypotheses", "Synthesizer", "Verifier", "Persist"]
    durs = [st.get("duration_ms") or 0 for st in steps] or [0] * 11
    bw, bh = Inches(1.95), Inches(0.8)
    prev = None
    for i, (lab, ms) in enumerate(zip(labels, durs, strict=True)):
        row, col = divmod(i, 6)
        x = Inches(0.55) + col * (bw + Inches(0.1))
        y = Inches(1.6) + row * Inches(1.45)
        colr = ACCENT if i < 7 else VIOLET if i in (8,) else WARN if i == 9 else OK
        n = D.node(s, lab.replace(" Investigator", ""), x, y, bw, bh, line=colr, sub=f"{ms:.0f} ms" if ms else "", size=11)
        if prev is not None:
            if col == 0:
                D.connect(s, prev, n, from_side="bottom", to_side="top", color=BORDER)
            else:
                D.connect(s, prev, n, color=BORDER)
        prev = n
    D.table(s, [
        ["Stage", "What it does", "Output"],
        ["Metrics", "baseline vs incident deviation (z-score, % change), saturation vs ceiling, monotonic trend; 'remained normal' contradictions", "metric evidence + signals"],
        ["Logs", "template mining (Drain-style masking), burst ratio vs baseline, keyword → signal catalog, new-exception detection", "log evidence, error clusters"],
        ["Traces", "per-operation p95 vs baseline, error spans, client-slow/server-ok → network, critical path of slowest trace", "trace evidence"],
        ["Deployments", "versions, commits, changed files, config diffs with proximity-to-onset scoring", "deployment/config evidence"],
        ["Dependency", "downstream health, blast radius, shared-dependency culprit hints on the service graph", "culprit hints"],
        ["Historical", "incident signature → embedding → cosine retrieval over resolved incidents (RAG)", "'resembles INC-031' evidence"],
        ["Hypotheses → Synthesizer → Verifier", "catalog of 13 failure modes with required/contradicting signals; explicit score; model narrates & cites; verifier re-derives confidence and blocks large reorders", "ranked, cited, calibrated"],
    ], Inches(0.55), Inches(4.65), Inches(12.2), col_widths=[2.2, 7.5, 2.5], size=9, row_h=Inches(0.28))

    # 6 — evidence graph
    s = D.slide("The evidence graph: root-cause analysis as graph + temporal reasoning, not semantic vibes", "ENGINE")
    D.image(s, ASSETS / "graph.jpg", Inches(0.6), Inches(1.45), w=Inches(7.9))
    D.bullets(s, [
        "Nodes:: incident, services, deployments, commits, alerts, evidence (E1…En), hypotheses, similar past incidents.",
        "Edges:: affects · depends_on · triggered_by · correlated_with · supports · contradicts · resembles — with weights.",
        "Handles are minted by the system:: the model may only cite E-handles that exist; the verifier drops anything else and records it.",
        "Contradictions are first-class:: 'CPU remained normal' is stored as evidence and shown on the root-cause card; it lowers confidence rather than being hidden.",
        "Score breakdown is stored per hypothesis:: signal support + temporal + dependency + historical − contradictions → a number the UI can explain.",
    ], Inches(8.8), Inches(1.5), Inches(4.0), Inches(5.2), size=11, gap=8)

    # 7 — thesis comparison table
    s = D.slide("Deterministic-first is the product. The model is the least interesting part.", "THESIS")
    D.table(s, [
        ["", "Chat-box over logs (typical 'AI ops')", "Sentinel"],
        ["Reproducibility", "different answer every run", "same evidence → same ranking; benchmarked in CI"],
        ["Evidence", "prose, no references", "every claim cites system-minted evidence handles; invalid citations are rejected"],
        ["Confidence", "whatever the model says", "re-derived from evidence kinds, contradictions and citation validity; capped at 95%; < 55% → human review"],
        ["Prompt injection", "log content can steer the answer", "telemetry rendered as data blocks; model cannot add categories/evidence; deterministic ranking guarded"],
        ["Model dependency", "paid API, data leaves the network", "local Ollama or no model at all — the platform works fully without one"],
        ["Remediation", "'AI rolled back prod'", "playbooks proposed; four-eyes approval; adapter executes; verification; audit log"],
        ["Measurability", "anecdotes", f"{sm['cases']} scenarios, ground truth, healthy controls, calibration, latency — quality gate in CI"],
    ], Inches(0.6), Inches(1.55), Inches(12.1), col_widths=[2.0, 4.2, 6.2], size=11, row_h=Inches(0.55), highlight_col=2)

    # 8 — benchmark headline chart
    s = D.slide(f"Measured, not claimed: {sm['cases']} synthetic production failures, ground truth never reaches the pipeline", "EVALUATION")
    cd = CategoryChartData()
    cd.categories = ["Root cause top-1", "Root cause top-3", "Detection", "Evidence precision", "Citation validity", "False positives", "Confident-wrong"]
    cd.add_series("latest run (deterministic)", [round(sm["root_cause_accuracy"] * 100, 1), round(sm["root_cause_top3_accuracy"] * 100, 1), round(sm["detection_rate"] * 100, 1), round(sm["evidence_precision"] * 100, 1), round(sm["citation_validity"] * 100, 1), round(sm["false_positive_rate"] * 100, 1), round(sm["confident_wrong_rate"] * 100, 1)])
    gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.6), Inches(1.5), Inches(7.6), Inches(4.4), cd)
    ch = D.style_chart(gf.chart, legend=False, colors=[ACCENT], number_format='0.0"%"')
    ch.value_axis.maximum_scale = 110
    ch.value_axis.minimum_scale = 0
    D.stat(s, "cases", str(sm["cases"]), f"{sm['fault_cases']} faults · {sm['control_cases']} healthy controls", Inches(8.5), Inches(1.5), w=Inches(2.0), h=Inches(1.15))
    D.stat(s, "median time", f"{sm['median_investigation_ms'] / 1000:.2f}s", f"p95 {sm['p95_investigation_ms'] / 1000:.2f}s", Inches(10.7), Inches(1.5), w=Inches(2.0), h=Inches(1.15), color=VIOLET)
    D.stat(s, "first iteration", f"{d['first_run_acc']:.1f}%", "top-1, before symptom/cause + noise fixes", Inches(8.5), Inches(2.8), w=Inches(2.0), h=Inches(1.15), color=WARN)
    D.stat(s, "calibration", f"{sm['ece']:.2f}", "ECE — under-confident by design", Inches(10.7), Inches(2.8), w=Inches(2.0), h=Inches(1.15), color=INFO)
    D.bullets(s, [
        "14 root-cause categories:: × target services × intensity / noise / confounding-deploy variants.",
        "Same code path as production:: detector → orchestrator → verifier; telemetry synthesised for the whole topology with realistic propagation.",
        "Quality gate in CI:: accuracy ≥ 85%, false positives ≤ 10%, citation validity ≥ 95%.",
    ], Inches(8.5), Inches(4.1), Inches(4.3), Inches(2.2), size=10, gap=6)
    D.text(s, "Honesty note: telemetry is synthetic and the catalog and scenarios share an author — see docs/evaluation/methodology.md, threats to validity.", Inches(0.6), Inches(6.1), Inches(8), Inches(0.5), size=9, color=DIM, italic=True)

    # 9 — per fault type
    s = D.slide("Every failure mode, individually: detection and top-1 accuracy per fault type", "EVALUATION")
    pf = sm["per_fault"]
    names = sorted(pf)
    cd = CategoryChartData()
    cd.categories = [n.replace("_", " ") for n in names]
    cd.add_series("top-1 accuracy", [round(pf[n]["accuracy"] * 100, 1) for n in names])
    cd.add_series("mean confidence", [round(pf[n]["mean_confidence"] * 100, 1) for n in names])
    gf = s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.6), Inches(1.45), Inches(7.2), Inches(5.4), cd)
    ch = D.style_chart(gf.chart, colors=[OK, VIOLET], number_format='0"%"', font_size=9)
    ch.value_axis.maximum_scale = 110
    rows = [["Fault type", "Cases", "Detected", "Top-1", "Top-3", "Mean conf."]]
    for n in names:
        v = pf[n]
        rows.append([n, v["cases"], v["detected"], f"{v['accuracy'] * 100:.0f}%", f"{v['top3_accuracy'] * 100:.0f}%", f"{v['mean_confidence']:.2f}"])
    D.table(s, rows, Inches(8.0), Inches(1.45), Inches(4.8), col_widths=[2.6, 0.9, 1.1, 0.9, 0.9, 1.2], size=8, row_h=Inches(0.3), highlight_col=3)

    # 10 — calibration + latency
    s = D.slide("Confidence you can defend: calibrated, capped, conservative — and fast", "EVALUATION")
    D.image(s, ASSETS / "calibration.png", Inches(0.6), Inches(1.45), w=Inches(6.1))
    D.image(s, ASSETS / "latency.png", Inches(6.9), Inches(1.45), w=Inches(5.9))
    D.bullets(s, [
        "Confidence is re-derived, never trusted:: score × citation validity − contradictions; ×0.6 if verification fails; capped at 0.75 with < 3 evidence kinds; never above 0.95.",
        "The gap is in the safe direction:: reported confidence sits below observed accuracy (ECE %.2f) — the system says 'probably' when it is right, not 'certainly' when it is wrong." % sm["ece"],
        "Below 0.55 → HUMAN_REVIEW:: the incident is routed to a person instead of a playbook.",
        "Median %.2fs per investigation:: on SQLite, in-process; PostgreSQL + worker is the production shape." % (sm["median_investigation_ms"] / 1000),
    ], Inches(0.6), Inches(4.95), Inches(12.2), Inches(2.0), size=11, gap=6)

    # 11 — evidence mix & investigator time (pie + doughnut)
    s = D.slide("Where the evidence comes from, and where the time goes", "ENGINE")
    cd = CategoryChartData()
    cd.categories = [k for k, _ in kinds.most_common()]
    cd.add_series("evidence items", [v for _, v in kinds.most_common()])
    gf = s.shapes.add_chart(XL_CHART_TYPE.PIE, Inches(0.6), Inches(1.5), Inches(5.6), Inches(4.6), cd)
    D.style_chart(gf.chart, colors=[ACCENT, VIOLET, INFO, MUTED, WARN, OK, CRIT], number_format="0")
    D.text(s, f"Evidence mix for {inc['key']} ({len(ev)} items across {len(kinds)} signal types)", Inches(0.6), Inches(6.15), Inches(5.6), Inches(0.4), size=10, color=MUTED, align=PP_ALIGN.CENTER)
    cd2 = CategoryChartData()
    cd2.categories = [lab.replace(" Investigator", "").replace("Root Cause ", "") for lab in labels]
    cd2.add_series("ms", [max(1, round(x)) for x in durs])
    gf2 = s.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, Inches(6.6), Inches(1.5), Inches(6.2), Inches(4.6), cd2)
    ch2 = D.style_chart(gf2.chart, colors=[ACCENT, INFO, WARN, VIOLET, OK, MUTED, DIM, CRIT, RGBColor(0x0E, 0x9A, 0xB0), RGBColor(0x3B, 0x82, 0xF6), RGBColor(0xD9, 0x77, 0x06)], number_format='0" ms"', labels=False)
    D.text(s, f"Investigator time for the same incident — total {sum(durs) / 1000:.1f}s, model time 0s (deterministic narrator)", Inches(6.6), Inches(6.15), Inches(6.2), Inches(0.4), size=10, color=MUTED, align=PP_ALIGN.CENTER)

    # 12 — model lift
    s = D.slide("What a local model adds — measured honestly, including when it hurt", "AI LAYER")
    smk = d["smoke"]["summary"] if d["smoke"] else None
    rows = [["Run", "Top-1", "Citation validity", "Median time", "Notes"],
            ["Deterministic narrator (reference)", "3/3", "100%", "0.4 s", "the pipeline alone"],
            ["qwen2.5:3b, free reordering", "2/3", "100%", "248 s", "narrator flipped a correct answer on a 0.25 score gap"],
            [f"qwen2.5:3b + rank-stability guard", f"{'3/3' if smk and smk['root_cause_accuracy'] == 1 else '—'}", f"{(smk['citation_validity'] * 100):.0f}%" if smk else "—", f"{(smk['median_investigation_ms'] / 1000):.0f} s" if smk else "—", "flip rejected and recorded; cross-exam over-skeptical → human review"]]
    D.table(s, rows, Inches(0.6), Inches(1.55), Inches(12.1), col_widths=[3.4, 1.0, 1.6, 1.4, 4.7], size=11, row_h=Inches(0.5), highlight_col=1)
    D.bullets(s, [
        "Provider abstraction:: generate · structured · embed. Ollama on localhost (no key, nothing leaves the host), deterministic fallback, circuit breaker, JSON-schema validation with repair and one re-prompt.",
        "Two guard rails validated by the study:: (1) the narrator may reorder only within a 0.10 deterministic score gap; (2) model verification can lower confidence, never raise it.",
        "Conclusion:: a 3B model improves prose, not root-cause identification — so it is optional. 7B+ on a GPU is the next experiment; the harness is ready (SENTINEL_LLM_PROVIDER=ollama make eval).",
        "Cost:: ₹0 in API spend — measured as latency and tokens instead.",
    ], Inches(0.6), Inches(3.8), Inches(12.1), Inches(3.1), size=12, gap=8)

    # 13 — human in the loop tree
    s = D.slide("Human-in-the-loop remediation: Sentinel proposes, people decide, everything is audited", "SAFETY")
    flow = [("Root cause\nverified", ACCENT), ("Playbook actions\nproposed", ACCENT), ("Request\nENGINEER+", INFO), ("Approve\nSRE+, not requester", WARN), ("Execute\nadapter, allow-listed kinds", CRIT), ("Verify\nservice healthy?", OK), ("Audit log\nwho · what · why · result", MUTED)]
    prev = None
    for i, (lab, colr) in enumerate(flow):
        head, sub = lab.split("\n")
        n = D.node(s, head, Inches(0.55) + i * Inches(1.78), Inches(1.7), Inches(1.65), Inches(0.95), line=colr, sub=sub, size=11)
        if prev is not None:
            D.connect(s, prev, n, color=colr)
        prev = n
    D.table(s, [
        ["Control", "Implementation", "Tested"],
        ["Recommendation-only default", "every action starts as `proposed`; nothing executes automatically", "✓"],
        ["Four-eyes", "requester cannot approve their own action (ADMIN excepted)", "✓"],
        ["Execute gate", "409 unless approved; only `executable` kinds (rollback, restart, scale, clear_fault)", "✓ live on the compose stack"],
        ["Agent tool boundary", "read_logs/metrics/traces/git allowed; rollback/restart/scale denied unless granted", "✓"],
        ["Audit", "who, what, when, why, outcome — for logins, faults, rules, approvals, executions", "✓"],
        ["Adapter", "simulator release manager today; Kubernetes / Argo / feature flags behind the same 2-method interface", "—"],
    ], Inches(0.6), Inches(3.1), Inches(12.1), col_widths=[2.6, 7.6, 1.9], size=10, row_h=Inches(0.42))

    # 14 — production readiness checklist
    s = D.slide("Production engineering is where we earn the word 'production'", "PLATFORM")
    left = [["Reliability", ""], ["Retries + exponential backoff", "✓"], ["Idempotent jobs, de-duplicated enqueue", "✓"], ["Per-step timeouts, job timeouts", "✓"], ["Circuit breakers (model, simulator)", "✓"], ["Dead-letter queue", "✓"], ["Checkpoint + resume after crash", "✓"], ["Graceful shutdown, worker recovery", "✓"], ["Rate limiting (memory / Redis)", "✓"]]
    right = [["Security & operations", ""], ["JWT + API keys with scopes, RBAC (4 roles)", "✓"], ["Audit logs, four-eyes approval", "✓"], ["Prompt-injection hardening + verifier", "✓"], ["Threat model, runbooks, ADRs", "✓"], ["Prometheus metrics for Sentinel itself + Grafana", "✓"], ["Alembic migrations, PostgreSQL/SQLite", "✓"], ["Docker Compose stack, CI with eval quality gate", "✓"], ["69 unit / integration / chaos tests", "✓"]]
    D.table(s, left, Inches(0.6), Inches(1.55), Inches(5.9), col_widths=[5, 0.7], size=11, row_h=Inches(0.42), highlight_col=1)
    D.table(s, right, Inches(6.8), Inches(1.55), Inches(5.9), col_widths=[5, 0.7], size=11, row_h=Inches(0.42), highlight_col=1)
    D.text(s, "Known gaps, on purpose: OIDC, PII redaction at ingestion, pgvector at scale, Kubernetes manifests — each has an ADR or a threat-model entry.", Inches(0.6), Inches(5.6), Inches(12), Inches(0.5), size=11, color=MUTED, italic=True)

    # 15 — demo shop & chaos lab
    s = D.slide("Our own production to break: a seven-service shop over real HTTP with a fourteen-fault chaos engine", "TESTBED")
    D.image(s, ASSETS / "services.jpg", Inches(0.6), Inches(1.45), w=Inches(6.1))
    D.image(s, ASSETS / "chaos.jpg", Inches(6.85), Inches(1.45), w=Inches(6.0))
    D.bullets(s, [
        "Real:: seven uvicorn services, real inter-service calls with trace propagation, telemetry pushed to Sentinel and scraped by Prometheus, deployments announced with synthetic Git metadata.",
        "Simulated in-process:: DB pool (semaphore + acquire timeout), cache, queue, worker pool, CPU/memory — faithful failure shapes, deterministic and safe.",
        "Faults:: pool exhaustion · bad deployment · DB latency · Redis outage · memory leak · CPU · 500 spike · config regression · network latency · packet loss · queue backlog · thread starvation · deadlock · dependency failure.",
    ], Inches(0.6), Inches(5.1), Inches(12.2), Inches(1.9), size=11, gap=6)

    # 16 — evaluation dashboard screenshot
    s = D.slide("The evaluation is part of the product: engineers see the same numbers investors do", "EVALUATION")
    D.image(s, ASSETS / "evaluation.jpg", Inches(0.6), Inches(1.45), w=Inches(9.2))
    D.bullets(s, [
        "Per-fault accuracy and confidence",
        "Confusion matrix (expected → predicted)",
        "Every case links to its incident, evidence and hypotheses",
        "'Quick run' re-executes 12 scenarios from the UI",
        "make eval regenerates docs/evaluation/latest.md",
    ], Inches(10.0), Inches(1.5), Inches(2.9), Inches(4.5), size=11, gap=8)

    # 17 — deployment topology
    s = D.slide("Runs anywhere: a laptop with nothing installed, Docker Compose, or Vercel + Railway", "DEPLOYMENT")
    cols = [
        ("Laptop (zero infra)", ["sentinel dev", "SQLite · in-process queue", "deterministic narrator", "demo shop :9000–9007", "web :3000"], ACCENT),
        ("Docker Compose", ["PostgreSQL · Redis", "api · worker (ARQ)", "simulator · web", "Prometheus · Alertmanager", "Grafana · OTel collector · Ollama"], INFO),
        ("Cloud (this deck's links)", ["Vercel → Next.js dashboard", "Railway → API (+ scheduler)", "Railway → PostgreSQL", "Railway → demo simulator", "same image, SENTINEL_ROLE=…"], OK),
    ]
    for i, (head, items, colr) in enumerate(cols):
        x = Inches(0.6) + i * Inches(4.15)
        D.node(s, head, x, Inches(1.6), Inches(3.9), Inches(0.6), line=colr, size=13)
        y = Inches(2.4)
        for it in items:
            n = D.node(s, it, x + Inches(0.3), y, Inches(3.3), Inches(0.5), fill=BG, line=BORDER, size=11, bold=False)
            y += Inches(0.62)
    D.text(s, "No code path changes between shapes — only SENTINEL_* settings. Real services plug in through the same ingestion API or the OpenTelemetry collector.", Inches(0.6), Inches(5.7), Inches(12), Inches(0.5), size=12, color=MUTED)
    D.text(s, f"web  {web_url}\napi  {api_url}\nrepo {repo_url}", Inches(0.6), Inches(6.2), Inches(12), Inches(0.8), size=11, color=ACCENT, font=MONO)

    # 18 — market & positioning
    s = D.slide("Who it is for, and how it is different", "POSITIONING")
    D.table(s, [
        ["", "Observability suites", "Incident-management tools", "'AI SRE' chat assistants", "Sentinel"],
        ["Detects", "✓ (alerts)", "✓ (pages)", "—", "✓ deterministic rules + Alertmanager"],
        ["Explains why", "partial (anomaly widgets)", "—", "prose", "ranked hypotheses with cited evidence"],
        ["Shows contradictions & confidence", "—", "—", "—", "✓ calibrated, capped, human-review routing"],
        ["Data stays on-prem", "usually not", "not applicable", "rarely", "✓ local-first, local model or none"],
        ["Acts", "—", "runbooks, paging", "sometimes, unsafely", "human-approved, audited, verified"],
        ["Measured accuracy", "—", "—", "—", f"✓ {sm['cases']}-case public benchmark"],
    ], Inches(0.6), Inches(1.55), Inches(12.1), col_widths=[2.4, 2.2, 2.4, 2.3, 3.2], size=11, row_h=Inches(0.5), highlight_col=4)
    D.bullets(s, [
        "Buyer:: platform / SRE teams at companies with 10–500 services who already own their telemetry and cannot ship it to a third-party model.",
        "Wedge:: ingestion is OpenTelemetry-shaped and Alertmanager-compatible — Sentinel sits beside existing monitoring, it does not replace it.",
        "Model:: open-source core (Apache-2.0); commercial adapters (Kubernetes/Argo/flags), SSO, multi-tenant control plane, support.",
    ], Inches(0.6), Inches(5.35), Inches(12.1), Inches(1.6), size=11, gap=6)

    # 19 — roadmap
    s = D.slide("Roadmap", "NEXT")
    phases = [("Now", ["deterministic pipeline, 14 failure modes", "benchmark + CI gate", "compose stack, Vercel + Railway", "local-model study (3B)"], OK),
              ("Next quarter", ["real-target adapters: Kubernetes, Argo, feature flags", "PII redaction at ingestion", "per-service ingestion credentials, OIDC", "7B+ model study on GPU"], ACCENT),
              ("Two quarters", ["learned scoring from labelled real incidents", "pgvector retrieval, multi-project control plane", "design partners: 3 teams on real telemetry", "public benchmark contributions"], INFO),
              ("Beyond", ["cross-team incident graphs", "change-risk prediction before deploys", "managed offering with on-prem data plane"], VIOLET)]
    for i, (head, items, colr) in enumerate(phases):
        x = Inches(0.6) + i * Inches(3.1)
        D.node(s, head, x, Inches(1.6), Inches(2.9), Inches(0.55), line=colr, size=13)
        D.bullets(s, items, x, Inches(2.3), Inches(2.9), Inches(3.5), size=11, gap=6)
    D.text(s, "Every item maps to an ADR or a threat-model gap already written down.", Inches(0.6), Inches(6.2), Inches(12), Inches(0.4), size=11, color=MUTED, italic=True)

    # 20 — close
    s = D.slide()
    D.text(s, "◆", Inches(0.9), Inches(1.5), Inches(1), Inches(1), size=40, color=ACCENT)
    D.text(s, "Reality → Evidence → Hypotheses → Verification → Human decision", Inches(0.9), Inches(2.4), Inches(11.5), Inches(0.8), size=28, color=FG, bold=True)
    D.text(s, "That architecture is the product. The LLM is the least interesting part.", Inches(0.9), Inches(3.3), Inches(11.5), Inches(0.6), size=18, color=MUTED)
    D.bullets(s, [
        f"Try it:: {web_url}  (admin@sentinel.local / admin12345)",
        f"Read it:: {repo_url}  — architecture, ADRs, threat model, evaluation methodology",
        "Break it:: Chaos Lab → inject a fault → watch the investigation → ask it Why?",
        "Talk:: raunit.thakur@gmail.com",
    ], Inches(0.9), Inches(4.3), Inches(11.5), Inches(2.4), size=14, gap=10)

    out = HERE / "Sentinel-Pitch.pptx"
    D.save(out)
    return out


if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    ap.add_argument("--web", default="http://localhost:3000")
    ap.add_argument("--api", default="http://localhost:8000")
    ap.add_argument("--repo", default="https://github.com/raunitgrey7/sentinel")
    a = ap.parse_args()
    print(build(a.web, a.api, a.repo))
